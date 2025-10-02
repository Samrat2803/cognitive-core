"""
PowerPoint Deck Generator

Compiles series of infographics into professional PowerPoint presentations.
Supports multiple layouts, section dividers, and export to PPTX/PDF.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import os
from typing import Dict, List, Any, Optional
from datetime import datetime
import json


class DeckGenerator:
    """Compile infographics into PowerPoint presentations"""
    
    # Standard slide dimensions (16:9)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(5.625)
    
    # Aistra color palette
    COLORS = {
        "primary": RGBColor(217, 243, 120),    # #d9f378
        "secondary": RGBColor(93, 83, 92),     # #5d535c
        "dark": RGBColor(51, 51, 51),          # #333333
        "darkest": RGBColor(28, 30, 32),       # #1c1e20
        "white": RGBColor(255, 255, 255),
        "light_gray": RGBColor(245, 245, 245),
    }
    
    def __init__(self, output_dir: str = "artifacts"):
        """
        Initialize deck generator
        
        Args:
            output_dir: Directory to save generated presentations
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def create_deck(
        self,
        infographic_paths: List[str],
        deck_config: Dict[str, Any],
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create PowerPoint from multiple infographics
        
        Args:
            infographic_paths: List of infographic image paths
            deck_config: Configuration
                {
                    "title": "Political Analysis Report",
                    "subtitle": "Q4 2025",
                    "author": "Analyst Team",
                    "date": "October 2, 2025",
                    "add_title_slide": True,
                    "add_section_dividers": True,
                    "sections": [
                        {
                            "title": "Sentiment Analysis",
                            "slides": [0, 1, 2],  # Indices in infographic_paths
                            "notes": "Key findings..."
                        },
                        {
                            "title": "Bias Detection",
                            "slides": [3, 4]
                        }
                    ],
                    "add_summary_slide": True,
                    "theme": "dark"  # dark or light
                }
            output_path: Custom output path
        
        Returns:
            Dict with artifact info
        """
        print(f"📊 Creating PowerPoint deck with {len(infographic_paths)} infographics...")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        
        # Get config
        title = deck_config.get("title", "Political Analysis Report")
        subtitle = deck_config.get("subtitle", "")
        author = deck_config.get("author", "Analysis Team")
        date = deck_config.get("date", datetime.now().strftime("%B %d, %Y"))
        add_title = deck_config.get("add_title_slide", True)
        add_dividers = deck_config.get("add_section_dividers", True)
        add_summary = deck_config.get("add_summary_slide", True)
        sections = deck_config.get("sections", [])
        theme = deck_config.get("theme", "dark")
        
        slide_count = 0
        
        # Title slide
        if add_title:
            print(f"   📄 Creating title slide...")
            self._create_title_slide(prs, title, subtitle, author, date, theme)
            slide_count += 1
        
        # Process sections
        if sections:
            for section in sections:
                section_title = section.get("title", "Section")
                slide_indices = section.get("slides", [])
                section_notes = section.get("notes", "")
                
                print(f"   📑 Section: {section_title} ({len(slide_indices)} slides)")
                
                # Section divider
                if add_dividers:
                    self._create_section_divider(prs, section_title, theme)
                    slide_count += 1
                
                # Add infographic slides
                for idx in slide_indices:
                    if 0 <= idx < len(infographic_paths):
                        img_path = infographic_paths[idx]
                        if os.path.exists(img_path):
                            self._create_content_slide(prs, img_path, theme, notes=section_notes)
                            slide_count += 1
                        else:
                            print(f"   ⚠️  Infographic not found: {img_path}")
        else:
            # No sections defined, just add all infographics
            print(f"   📄 Adding all infographics as slides...")
            for img_path in infographic_paths:
                if os.path.exists(img_path):
                    self._create_content_slide(prs, img_path, theme)
                    slide_count += 1
                else:
                    print(f"   ⚠️  Infographic not found: {img_path}")
        
        # Summary slide
        if add_summary:
            print(f"   📄 Creating summary slide...")
            self._create_summary_slide(prs, deck_config, theme)
            slide_count += 1
        
        # Generate output path
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"deck_{timestamp}.pptx"
            output_path = os.path.join(self.output_dir, filename)
        
        # Save presentation
        print(f"   💾 Saving to: {output_path}")
        prs.save(output_path)
        
        print(f"   ✅ Deck created with {slide_count} slides!")
        
        # Return artifact info
        return {
            "artifact_id": f"deck_{datetime.now().timestamp()}",
            "type": "powerpoint_deck",
            "title": title,
            "path": output_path,
            "num_slides": slide_count,
            "num_sections": len(sections) if sections else 0,
            "theme": theme,
            "format": "PPTX",
            "size_bytes": os.path.getsize(output_path)
        }
    
    def _create_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: str,
        author: str,
        date: str,
        theme: str
    ):
        """Create title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["darkest"] if theme == "dark" else self.COLORS["white"]
        
        # Top accent bar
        bar = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            self.SLIDE_WIDTH,
            Inches(0.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS["primary"]
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.2),
                Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            
            p = subtitle_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.color.rgb = self.COLORS["white"] if theme == "dark" else self.COLORS["dark"]
            p.alignment = PP_ALIGN.CENTER
        
        # Footer (author + date)
        footer_text = f"{author} | {date}"
        footer_box = slide.shapes.add_textbox(
            Inches(1), Inches(5),
            Inches(8), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = footer_text
        
        p = footer_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS["secondary"]
        p.alignment = PP_ALIGN.CENTER
    
    def _create_section_divider(
        self,
        prs: Presentation,
        section_title: str,
        theme: str
    ):
        """Create section divider slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["dark"] if theme == "dark" else self.COLORS["light_gray"]
        
        # Side accent bar (vertical)
        bar = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            Inches(0.3),
            self.SLIDE_HEIGHT
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS["primary"]
        bar.line.fill.background()
        
        # Section title (centered)
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2),
            Inches(7), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["white"] if theme == "dark" else self.COLORS["darkest"]
        p.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(
        self,
        prs: Presentation,
        infographic_path: str,
        theme: str,
        notes: str = ""
    ):
        """Create slide with infographic"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["darkest"] if theme == "dark" else self.COLORS["white"]
        
        # Get image dimensions
        try:
            img = Image.open(infographic_path)
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            
            # Calculate size to fit slide (leave margins)
            max_width = Inches(9)
            max_height = Inches(5)
            
            if img_aspect > (max_width / max_height):
                # Width is limiting factor
                pic_width = max_width
                pic_height = pic_width / img_aspect
            else:
                # Height is limiting factor
                pic_height = max_height
                pic_width = pic_height * img_aspect
            
            # Center image
            left = (self.SLIDE_WIDTH - pic_width) / 2
            top = (self.SLIDE_HEIGHT - pic_height) / 2
            
            # Add image
            slide.shapes.add_picture(
                infographic_path,
                left, top,
                width=pic_width,
                height=pic_height
            )
            
            # Add notes if provided
            if notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes
        
        except Exception as e:
            print(f"   ⚠️  Could not add image {infographic_path}: {e}")
    
    def _create_summary_slide(
        self,
        prs: Presentation,
        deck_config: Dict,
        theme: str
    ):
        """Create summary/closing slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["darkest"] if theme == "dark" else self.COLORS["white"]
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Summary"
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER
        
        # Summary points
        summary_text = deck_config.get("summary", "Thank you for viewing this analysis.")
        
        text_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.8),
            Inches(6), Inches(2)
        )
        text_frame = text_box.text_frame
        text_frame.text = summary_text
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = self.COLORS["white"] if theme == "dark" else self.COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER
        
        # Contact/footer
        footer_text = deck_config.get("contact", "For more information, contact the analysis team.")
        footer_box = slide.shapes.add_textbox(
            Inches(2), Inches(5),
            Inches(6), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = footer_text
        
        p = footer_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = self.COLORS["secondary"]
        p.alignment = PP_ALIGN.CENTER
    
    def export_to_pdf(self, pptx_path: str, pdf_path: Optional[str] = None) -> str:
        """
        Export PowerPoint to PDF (requires LibreOffice or similar)
        
        Note: This is a placeholder. Actual PDF conversion requires:
        - LibreOffice: soffice --headless --convert-to pdf file.pptx
        - Or python-pptx-pdf library
        - Or comtypes on Windows
        
        Args:
            pptx_path: Path to PPTX file
            pdf_path: Output PDF path
        
        Returns:
            Path to PDF file
        """
        if pdf_path is None:
            pdf_path = pptx_path.replace(".pptx", ".pdf")
        
        print(f"📄 PDF export: {pdf_path}")
        print("   ⚠️  Note: PDF export requires LibreOffice or similar tool")
        print("   Run: soffice --headless --convert-to pdf " + pptx_path)
        
        return pdf_path


# Standalone test
if __name__ == "__main__":
    print("📊 Testing Deck Generator...")
    print()
    
    # Check for test infographics
    import glob
    infographic_files = sorted(glob.glob("artifacts/infographic_*.png"))
    
    if not infographic_files:
        print("⚠️  No test infographics found!")
        print("Please run infographic_generator.py first to create test images.")
        exit(1)
    
    print(f"📸 Found {len(infographic_files)} infographics to use")
    print()
    
    generator = DeckGenerator(output_dir="artifacts")
    
    # Test 1: Simple deck with all infographics
    print("Test 1: Basic deck with all infographics")
    print("-" * 50)
    
    config_1 = {
        "title": "Political Analysis Report",
        "subtitle": "Comprehensive Sentiment & Bias Analysis",
        "author": "AI Analysis Team",
        "date": "October 2, 2025",
        "add_title_slide": True,
        "add_section_dividers": False,
        "add_summary_slide": True,
        "theme": "dark",
        "summary": "This analysis covered sentiment patterns across multiple countries and platforms.",
        "contact": "Generated by Political Analyst AI | analyst.ai"
    }
    
    try:
        result_1 = generator.create_deck(
            infographic_paths=infographic_files[:3],  # Use first 3
            deck_config=config_1
        )
        
        print(f"\n✅ Test 1 complete!")
        print(f"   Path: {result_1['path']}")
        print(f"   Slides: {result_1['num_slides']}")
        print(f"   Size: {result_1['size_bytes'] / 1024:.1f} KB")
        print()
    except Exception as e:
        print(f"❌ Test 1 failed: {e}")
        import traceback
        traceback.print_exc()
        print()
    
    # Test 2: Sectioned deck
    if len(infographic_files) >= 6:
        print("Test 2: Sectioned deck with dividers")
        print("-" * 50)
        
        config_2 = {
            "title": "Advanced Political Analysis",
            "subtitle": "Multi-Section Deep Dive",
            "author": "Research Division",
            "date": "October 2025",
            "add_title_slide": True,
            "add_section_dividers": True,
            "add_summary_slide": True,
            "theme": "dark",
            "sections": [
                {
                    "title": "Sentiment Analysis",
                    "slides": [0, 1, 2],
                    "notes": "This section covers sentiment trends across regions."
                },
                {
                    "title": "Platform Comparison",
                    "slides": [3, 4],
                    "notes": "Comparing different social media platforms."
                },
                {
                    "title": "Conclusions",
                    "slides": [5],
                    "notes": "Final findings and recommendations."
                }
            ],
            "summary": "Key insights: Positive sentiment trends, regional variations observed.",
            "contact": "Contact: analysis@political-ai.com"
        }
        
        try:
            result_2 = generator.create_deck(
                infographic_paths=infographic_files[:6],
                deck_config=config_2
            )
            
            print(f"\n✅ Test 2 complete!")
            print(f"   Path: {result_2['path']}")
            print(f"   Slides: {result_2['num_slides']}")
            print(f"   Sections: {result_2['num_sections']}")
            print()
        except Exception as e:
            print(f"❌ Test 2 failed: {e}")
            import traceback
            traceback.print_exc()
            print()
    
    print("🎉 Deck generator testing complete!")
    print(f"📁 Check the artifacts/ folder for generated PowerPoint files")
    print()
    
    # Save test results
    print("💾 Saving test results summary...")
    summary = {
        "test_date": datetime.now().isoformat(),
        "infographics_tested": len(infographic_files),
        "tests_run": 2 if len(infographic_files) >= 6 else 1,
        "output_directory": "artifacts/"
    }
    
    with open("artifacts/deck_test_results.json", 'w') as f:
        json.dump(summary, f, indent=2)
    
    print("✅ Test summary saved to: artifacts/deck_test_results.json")

